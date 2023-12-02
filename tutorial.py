import pptx
import os
import requests
import pandas as pd
import json
from Bio import Entrez
from pptx.util import Pt
import sys
import matplotlib.pyplot as plt
from pptx.util import Inches

presentation = pptx.Presentation(os.path.join("powerpoints", "input.pptx"))

def list_text_boxes(presentation, slide_num):
    slide = presentation.slides[slide_num-1]
    text_boxes = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text:
            text_boxes.append(shape.text)
    return text_boxes

# for idx, text in enumerate(list_text_boxes(presentation, 2), 1):
#     print(f"Text Box {idx}: {text}")

uniprot_id = sys.argv[1]

def update_text_of_textbox(presentation, slide, text_box_id, new_text):
    slide = presentation.slides[(slide - 1)]
    count = 0
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text:
            count += 1
            if count == text_box_id:
                text_frame = shape.text_frame
                first_paragraph = text_frame.paragraphs[0]
                first_run = first_paragraph.runs[0] if first_paragraph.runs else first_paragraph.add_run()
                # Preserve formatting of the first run
                font = first_run.font
                font_name = font.name
                font_size = font.size
                font_bold = font.bold
                font_italic = font.italic
                font_underline = font.underline
                font_color = font.color.rgb
                # Clear existing text and apply new text with preserved formatting
                text_frame.clear()  # Clears all text and formatting
                new_run = text_frame.paragraphs[0].add_run()  # New run in first paragraph
                new_run.text = new_text
                # Reapply formatting
                new_run.font.name = font_name
                new_run.font.size = font_size
                new_run.font.bold = font_bold
                new_run.font.italic = font_italic
                new_run.font.underline = font_underline
                new_run.font.color.rgb = font_color
                return
            

# SLIDE 1
update_text_of_textbox(presentation, 1, 1, uniprot_id)


# SLIDE 2
# UniProt API endpoint
UNIPROT_API_ENDPOINT = "https://rest.uniprot.org/uniprotkb/{}"


def retrieve_data(api_endpoint, gene_id):
    """Retrieves data from the given API endpoint using the given gene ID.
    Args:
    api_endpoint: The API endpoint to retrieve data from.
    gene_id: The gene ID to use in the API query.
    Returns:
    A dictionary containing the retrieved data.
    """
    response = requests.get(api_endpoint.format(gene_id))
    response.raise_for_status()
    return response.json()


uniprot_data = retrieve_data(UNIPROT_API_ENDPOINT, uniprot_id)

gene_name = uniprot_data["genes"][0]["geneName"]["value"]

official_gene_name = [database["properties"][0]["value"] for database in uniprot_data["uniProtKBCrossReferences"] if database["database"] == "PANTHER"][0]

ensembl_id = [database["id"] for database in uniprot_data["uniProtKBCrossReferences"] if database["database"] == "OpenTargets"][0]

sequence_similarity = [comment["texts"][0]["value"] for comment in uniprot_data["comments"] if comment["commentType"] == "SIMILARITY"][0]


update_text_of_textbox(presentation, 2, 5, gene_name)
update_text_of_textbox(presentation, 2, 10, official_gene_name)
update_text_of_textbox(presentation, 2, 8, ensembl_id)
update_text_of_textbox(presentation, 2, 7, uniprot_id)
update_text_of_textbox(presentation, 2, 9, sequence_similarity)


# SLIDE 3
def get_gene_summary(gene_id):
    # Fetch gene details
    handle = Entrez.efetch(db="gene", id=gene_id, retmode="xml")
    record = Entrez.read(handle)
    # Extract summary information
    gene_summary = record[0]['Entrezgene_summary']
    return gene_summary

ncbi_id = [database["id"] for database in uniprot_data["uniProtKBCrossReferences"] if database["database"] == "GeneID" ]
function_summary_value = get_gene_summary(ncbi_id)
update_text_of_textbox(presentation, 3, 2, function_summary_value)


# SLIDE 5
protein_sequence = uniprot_data["sequence"]['value']
protein_length = str(uniprot_data["sequence"]['length'])
protein_molecular_weight = str(uniprot_data["sequence"]['molWeight'])

update_text_of_textbox(presentation, 5, 3, protein_sequence)
update_text_of_textbox(presentation, 5, 4, protein_length)
update_text_of_textbox(presentation, 5, 5, protein_molecular_weight)


# SLIDE 4
# Build query string to get general information about AR and genetic constraint and tractability assessments 
query_string = """
    query ExpressionTest {{
        target(ensemblId: "{0}") {{
            id
            approvedName
            approvedSymbol
            expressions{{
            tissue{{
                id
                label
                anatomicalSystems
                organs
            }}
            rna{{
                zscore
                value
                unit
                level
            }}
            protein{{
                reliability
                level
                cellType{{
                    reliability
                    name
                    level
                }}
            }}
            }}
        }}
    }}
""".format(ensembl_id)

# Set variables object of arguments to be passed to endpoint
open_targets_variables = {"ensemblId": ensembl_id}
# Set base URL of GraphQL API endpoint
base_url = "https://api.platform.opentargets.org/api/v4/graphql"
# Perform POST request and check status code of response
r = requests.post(base_url, json={"query": query_string, "variables": open_targets_variables})
# Transform API response from JSON into Python dictionary
api_response = json.loads(r.text)
# Extracting the list of expressions
expressions = api_response['data']['target']['expressions']
# Creating a dataframe
df = pd.DataFrame(expressions)
# Extracting relevant columns
df = df[['tissue', 'rna', 'protein']]
df['tissue_label'] = df['tissue'].apply(lambda x: x['label'])
df['organ'] = df['tissue'].apply(lambda x: x['organs'])
df['rna_value'] = df['rna'].apply(lambda x: x['value'])
df['rna_unit'] = df['rna'].apply(lambda x: x['unit'])
df['protein_reliability'] = df['protein'].apply(lambda x: x['reliability'])
df['protein_level'] = df['protein'].apply(lambda x: x['level'])
# Dropping original columns
df = df.drop(columns=['tissue', 'rna', 'protein'])

#@@@
top_10_rna = df.sort_values(by='rna_value', ascending=False).head(10)
top_10_protein = df.sort_values(by='protein_level', ascending=False).head(10)

# Adjust protein level by adding one to each value
top_10_protein['adjusted_protein_level'] = top_10_protein['protein_level'] + 1

# Plot for top 10 RNA values
plt.figure(figsize=(10, 8))
plt.barh(top_10_rna['tissue_label'], top_10_rna['rna_value'], color='skyblue')
plt.xlabel('RNA count')
plt.ylabel('Tissue')
plt.title('Top 10 RNA expression by Tissue')
plt.gca().invert_yaxis()
plt.tight_layout()
plt.savefig('scratch/rna_plot.png')  # Save as PNG file
plt.close()

# Plot for top 10 Protein Levels by Tissue
plt.figure(figsize=(10, 8))
plt.barh(top_10_protein['tissue_label'], top_10_protein['adjusted_protein_level'], color='salmon')
plt.xlabel('Protein Level (arbitrary units)')
plt.ylabel('Tissue')
plt.title('Top 10 Protein Levels by Tissue')
plt.gca().invert_yaxis()
plt.tight_layout()
plt.savefig('scratch/protein_plot.png')  # Save as PNG file
plt.close()
#@@@

def add_image_to_slide(slide, image_path, left, top, width, height):
    slide.shapes.add_picture(image_path, left, top, width, height)

# Add RNA plot to a specific slide (e.g., slide number 1)
add_image_to_slide(presentation.slides[3], 'scratch/rna_plot.png', Inches(11), Inches(3), Inches(8) , Inches(6))

# Add Protein plot to another slide (e.g., slide number 2)
add_image_to_slide(presentation.slides[3], 'scratch/protein_plot.png', Inches(1), Inches(3), Inches(8) , Inches(6))

# Save the updated presentation 
output_path = os.path.join("powerpoints", uniprot_id + "_SUMMARY.pptx")
presentation.save(output_path)


def main():
    # Check if there are any arguments (excluding the script name)
    if len(sys.argv) > 1:
        print(f"Argument received: {sys.argv[1]}")
    else:
        print("No arguments provided.")

if __name__ == "__main__":
    main()





# SOME EXTRA FUNCTIONS YOU MAY FIND USEFUL
##########################################
# def update_size_of_textbox(presentation, slide, text_box_id, font_size):
#     slide = presentation.slides[(slide - 1)]
#     count = 0
#     for shape in slide.shapes:
#         if shape.has_text_frame and shape.text:
#             count += 1
#             if count == text_box_id:
#                 # Loop through paragraphs and runs to change font size
#                 for paragraph in shape.text_frame.paragraphs:
#                     for run in paragraph.runs:
#                         run.font.size = Pt(font_size)
#                 return

# def make_text_bold(presentation, slide, text_box_id):
#     slide = presentation.slides[(slide - 1)]
#     count = 0
#     for shape in slide.shapes:
#         if shape.has_text_frame and shape.text:
#             count += 1
#             if count == text_box_id:
#                 # Loop through paragraphs and runs to change font to bold
#                 for paragraph in shape.text_frame.paragraphs:
#                     for run in paragraph.runs:
#                         run.font.bold = True
#                 return

# def add_hyperlink_to_textbox(presentation, slide, text_box_id, url):
#     slide = presentation.slides[(slide - 1)]
#     count = 0
#     for shape in slide.shapes:
#         if shape.has_text_frame and shape.text:
#             count += 1
#             if count == text_box_id:
#                 # Extracting the first paragraph of the shape
#                 paragraph = shape.text_frame.paragraphs[0]
#                 for run in paragraph.runs:
#                     # print(f"Processing run with text: {run.text}")
#                     run.hyperlink.address = url
#                     break
#                 return